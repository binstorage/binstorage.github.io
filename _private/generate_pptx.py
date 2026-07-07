#!/usr/bin/env python3
"""
이력서 & 경력기술서 PPT 생성 - Modern Resume 템플릿 스타일
실행: python3 _private/generate_pptx.py
의존성: pip3 install python-pptx

출력 파일:
  _private/이력서.pptx
  _private/경력기술서.pptx
  _private/이력서_경력기술서.pptx  (통합본)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

# ── 경로 ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "modern-resume-template.pptx")
OUT_DIR = SCRIPT_DIR

# ── Design Tokens ──
ORANGE = RGBColor(0xE4, 0x79, 0x3C)
DARK_ORANGE = RGBColor(0xC8, 0x51, 0x03)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
TAG_BG = RGBColor(0xFD, 0xF0, 0xE6)

FONT_HEADING = "Poppins Medium"
FONT_BODY = "Poppins Light"
SW = Inches(10)
SH = Inches(5.625)
MAX_Y = Inches(5.15)  # 슬라이드 하단 안전 마진
L = Inches(0.6)       # 공통 왼쪽 마진


# ═══════════════════════════════════════
#  유틸리티
# ═══════════════════════════════════════
def delete_all_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[8])


def new_content_slide(prs):
    """콘텐츠 슬라이드 생성 (오렌지 top bar + 시작 y 반환)"""
    s = add_blank_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    return s, Inches(0.25)


def add_orange_sidebar(slide, width=Inches(3.2)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def needs_new_slide(y, needed):
    """y + needed > MAX_Y 이면 True"""
    return (y + needed) > MAX_Y


def textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
            font_size=Pt(12), bold=False, color=BLACK, align=PP_ALIGN.LEFT,
            line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def multi_line_textbox(slide, left, top, width, height, text, font_size=Pt(11),
                       color=BLACK, line_spacing=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = font_size
    run.font.color.rgb = color
    return txBox


def section_title(slide, left, top, text, font_size=Pt(18)):
    """섹션 타이틀 + 텍스트 길이에 맞는 오렌지 언더라인"""
    textbox(slide, left, top, Inches(8), Inches(0.32), text,
            font_name=FONT_HEADING, font_size=font_size, bold=True, color=BLACK)
    # 언더라인 길이: 글자 수 × 대략 글자폭
    char_width = font_size * 0.65  # 대략적 글자폭 비율
    underline_w = max(len(text) * char_width, Inches(0.5))
    underline_w = min(underline_w, Inches(8))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top + Inches(0.33), underline_w, Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    return top + Inches(0.48)


def sub_section(slide, left, top, text, font_size=Pt(11)):
    textbox(slide, left, top, Inches(8), Inches(0.25), text,
            font_name=FONT_HEADING, font_size=font_size, bold=True, color=DARK_ORANGE)
    return top + Inches(0.27)


def bullet_list(slide, left, top, items, width=Inches(8.5), font_size=Pt(9), color=GRAY):
    height = Inches(len(items) * 0.2 + 0.05)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = Pt(15)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.name = FONT_BODY
        run.font.size = font_size
        run.font.color.rgb = color
    return height  # 실제 사용한 높이 반환


def bullet_height(items):
    """bullet list가 차지할 높이 추정"""
    return Inches(len(items) * 0.2 + 0.05)


def add_table(slide, left, top, headers, rows, col_widths, row_h=Inches(0.28)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, row_h * n_rows)
    tbl = shape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.name = FONT_HEADING
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_margin(cell)
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.name = FONT_BODY
            run.font.size = Pt(9)
            run.font.color.rgb = BLACK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
            else:
                cell.fill.background()
            _set_cell_margin(cell)
    _set_table_borders(tbl, LIGHT_GRAY)
    return n_rows * row_h  # 테이블 높이 반환


def skill_tags_row(slide, left, top, tags, width=Inches(8.5)):
    x = left
    y = top
    tag_h = Inches(0.25)
    gap = Inches(0.08)
    max_x = left + width
    for tag in tags:
        tag_w = Pt(len(tag) * 7.5 + 22)
        if x + tag_w > max_x:
            x = left
            y += Inches(0.32)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tag_w, tag_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = TAG_BG
        shape.line.color.rgb = ORANGE
        shape.line.width = Pt(0.75)
        shape.adjustments[0] = 0.25
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tag
        run.font.name = FONT_HEADING
        run.font.size = Pt(8)
        run.font.color.rgb = DARK_ORANGE
        run.font.bold = True
        x += tag_w + gap
    return y + Inches(0.32)


def tags_height(tags, width=Inches(8.5)):
    """skill tags가 차지할 높이 추정"""
    x = 0
    rows = 1
    for tag in tags:
        tag_w = Pt(len(tag) * 7.5 + 22) + Inches(0.08)
        if x + tag_w > width:
            x = 0
            rows += 1
        x += tag_w
    return Inches(rows * 0.32)


def timeline_header(slide, left, top, title, period, width=Inches(8.5)):
    textbox(slide, left, top, Inches(5.5), Inches(0.28), title,
            font_name=FONT_HEADING, font_size=Pt(12), bold=True, color=BLACK)
    textbox(slide, left + Inches(5.5), top, Inches(3.0), Inches(0.28), period,
            font_name=FONT_BODY, font_size=Pt(9), color=GRAY, align=PP_ALIGN.RIGHT)
    return top + Inches(0.3)


def role_text(slide, left, top, text):
    textbox(slide, left, top, Inches(8), Inches(0.22), text,
            font_name=FONT_BODY, font_size=Pt(9), color=GRAY)
    return top + Inches(0.24)


def thin_divider(slide, left, top, width=Inches(8.8)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()
    return top + Inches(0.15)


def _set_cell_margin(cell, m=Emu(45720)):
    cell.margin_left = m
    cell.margin_right = m
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)


def _set_table_borders(table, color):
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(border_name))
                if ln is None:
                    ln = etree.SubElement(tcPr, qn(border_name))
                ln.set('w', '6350')
                ln.set('cap', 'flat')
                ln.set('cmpd', 'sng')
                solidFill = ln.find(qn('a:solidFill'))
                if solidFill is not None:
                    ln.remove(solidFill)
                solidFill = etree.SubElement(ln, qn('a:solidFill'))
                srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgb.set('val', str(color))


def estimate_project_height(overview, tasks, achievements, tags):
    """프로젝트 블록이 차지할 높이 추정"""
    h = Inches(0.3 + 0.24)  # timeline_header + role_text
    if overview:
        h += Inches(0.3)
    if tasks:
        h += Inches(0.27) + bullet_height(tasks)
    if achievements:
        h += Inches(0.27) + bullet_height(achievements)
    if tags:
        h += Inches(0.27) + tags_height(tags)
    h += Inches(0.1)
    return h


def render_project(slide, y, title, period, role_str, overview, tasks, achievements, tags):
    """프로젝트 블록을 렌더링하고 새 y 반환"""
    y = timeline_header(slide, L, y, title, period)
    y = role_text(slide, L + Inches(0.1), y, role_str)

    if overview:
        multi_line_textbox(slide, L + Inches(0.2), y, Inches(8.4), Inches(0.25),
                          overview, font_size=Pt(9), color=GRAY, line_spacing=Pt(13))
        y += Inches(0.28)

    if tasks:
        y = sub_section(slide, L + Inches(0.1), y, "담당 역할", font_size=Pt(10))
        bh = bullet_list(slide, L + Inches(0.2), y, tasks, width=Inches(8.3), font_size=Pt(9))
        y += bh + Inches(0.03)

    if achievements:
        y = sub_section(slide, L + Inches(0.1), y, "성과", font_size=Pt(10))
        bh = bullet_list(slide, L + Inches(0.2), y, achievements, width=Inches(8.3), font_size=Pt(9))
        y += bh + Inches(0.03)

    if tags:
        y = sub_section(slide, L + Inches(0.1), y, "사용 기술", font_size=Pt(10))
        y = skill_tags_row(slide, L + Inches(0.2), y, tags, width=Inches(8.3))

    y += Inches(0.08)
    return y


# ═══════════════════════════════════════
#  이력서 슬라이드 생성
# ═══════════════════════════════════════
def build_resume_slides(prs):
    # ── Slide 1: 사이드바 + 자기소개 ──
    s = add_blank_slide(prs)
    add_orange_sidebar(s, width=Inches(3.2))
    textbox(s, Inches(0.4), Inches(1.0), Inches(2.6), Inches(0.5), "안광빈",
            font_name=FONT_HEADING, font_size=Pt(30), bold=True, color=WHITE)
    textbox(s, Inches(0.4), Inches(1.6), Inches(2.6), Inches(0.35), "macOS, iOS Developer",
            font_name=FONT_BODY, font_size=Pt(13), color=WHITE)
    textbox(s, Inches(0.4), Inches(2.15), Inches(2.6), Inches(0.15), "\u2500" * 20,
            font_name=FONT_BODY, font_size=Pt(7), color=RGBColor(0xFF, 0xAA, 0x77))
    for text, cy in [
        ("\u260E  010-3499-5636", Inches(2.5)),
        ("\u2709  rhkdqlsdml10@gmail.com", Inches(2.8)),
        ("\u2302  서울 동작구 신대방동", Inches(3.1)),
    ]:
        textbox(s, Inches(0.4), cy, Inches(2.6), Inches(0.25), text,
                font_name=FONT_BODY, font_size=Pt(9), color=WHITE)

    rx = Inches(3.6)
    y = Inches(0.35)
    y = section_title(s, rx, y, "자기소개", font_size=Pt(18))
    intro = (
        '"그냥 해봐 할 수 있어!" 가장 제가 좋아하는 말입니다.\n\n'
        "[네트워크 보안 및 크로스플랫폼 개발]\n"
        "현재 제로트러스트 기반 네트워크 보안 회사에서 iOS, macOS 환경의 보안 에이전트를 "
        "개발하고 있습니다. Swift, UIKit, RxSwift에서 시작해 SwiftUI, Combine, Extension, "
        "Daemon, XPC, VPN, 보안 등으로 영역을 넓혀왔습니다.\n\n"
        "[언어 간 호환성에 대한 깊이 있는 이해]\n"
        "C++ 기반 네트워크 모듈을 Swift에서 호출하기 위해 C++, Objective-C, Objective-C++의 "
        "호환성과 컴파일러/링킹 과정까지 깊이 있게 학습했습니다.\n\n"
        "[아키텍처 설계 및 모듈화]\n"
        "모듈화와 아키텍처 설계에 관심이 많아 Tuist, Clean Architecture, TCA를 "
        "사이드 프로젝트에서 먼저 적용 후 회사에도 도입했습니다.\n\n"
        "[성장성 및 대규모 서비스에 대한 관심]\n"
        "대규모 서비스에도 관심이 많아 카카오뱅크 등을 참고하여 RIBs, ReactorKit을 "
        "사이드 프로젝트에 적용하고 있습니다."
    )
    multi_line_textbox(s, rx, y, Inches(6.0), Inches(3.8), intro,
                       font_size=Pt(9), color=GRAY, line_spacing=Pt(14))

    # ── Slide 2: 경력 + 학력 + 병역 ──
    s2, y = new_content_slide(prs)
    y = section_title(s2, L, y, "경력 사항")
    th = add_table(s2, L, y,
              ["기간", "회사명", "직책/직급", "담당 업무"],
              [
                  ("2023.04 - 현재", "프라이빗테크놀로지(주)", "연구소 주임연구원",
                   "iOS/macOS 보안 에이전트 개발 (SASE)"),
                  ("2020.12 - 2022.02", "제이엠트루", "개발연구소 연구원",
                   "iOS 앱 개발 및 배포"),
              ],
              [Inches(2.0), Inches(2.5), Inches(2.0), Inches(2.7)])
    y += th + Inches(0.15)

    y = section_title(s2, L, y, "학력")
    th = add_table(s2, L, y,
              ["기간", "학교명", "전공", "비고"],
              [
                  ("2014.03 - 2021.03", "청주대학교", "컴퓨터정보공학과", "졸업 (3.7/4.5)"),
                  ("2011 - 2014", "청석고등학교", "-", "졸업"),
              ],
              [Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.0)])
    y += th + Inches(0.05)
    multi_line_textbox(s2, L + Inches(0.2), y, Inches(8.3), Inches(0.3),
        "졸업작품: 캡스톤 디자인 - 고성능 헬멧 및 애플리케이션 제작 (블루투스, 충격감지, GPS 모듈 연동 iOS 앱 개발)",
        font_size=Pt(8), color=GRAY, line_spacing=Pt(12))
    y += Inches(0.28)

    y = section_title(s2, L, y, "병역사항")
    add_table(s2, L, y,
              ["구분", "기간", "군별/계급", "비고"],
              [("군필", "2015.08 - 2017.05", "육군 병장", "제대")],
              [Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.0)])

    # ── Slide 3: 자격증 + 인턴 + 기술 스택 ──
    s3, y = new_content_slide(prs)
    y = section_title(s3, L, y, "자격증 / 수상")
    th = add_table(s3, L, y,
              ["취득일", "자격증/수상명", "발급기관"],
              [("2020년", "공공데이터 활용 제품 및 서비스 장려상", "충북과학기술혁신원")],
              [Inches(2.0), Inches(4.2), Inches(3.0)])
    y += th + Inches(0.05)
    multi_line_textbox(s3, L + Inches(0.2), y, Inches(8.3), Inches(0.4),
        "착한가격 업소 정보를 공공데이터 API로 제공하는 서비스 앱 개발 (iOS 1명, 디자인 1명). "
        "iOS 개발 및 발표 담당. Alamofire를 활용한 공공데이터 API 통신, 내비게이션·리뷰 기능 구현.",
        font_size=Pt(8), color=GRAY, line_spacing=Pt(12))
    y += Inches(0.35)

    y = section_title(s3, L, y, "인턴 / 대외활동")
    y = timeline_header(s3, L + Inches(0.1), y, "LAB301", "2018.12 - 2020.07 (1년 8개월)")
    multi_line_textbox(s3, L + Inches(0.2), y, Inches(8.3), Inches(0.7),
        "학부시절 같은 학과 학부생 5명(Android 2, Backend 1, Design 1, iOS 1)이 팀을 꾸리고 "
        "사업자 등록 후 외주/팀프로젝트를 진행. 크몽, 오투잡 등 프리랜서 업체에 등록하여 "
        "앱스토어 배포, 구동 영상 프로토타입 등 약 5개 프로젝트 수행. "
        "외주 프로젝트 중 2개를 앱스토어 배포하여 사용자 100명+ 가입 달성.",
        font_size=Pt(9), color=GRAY, line_spacing=Pt(13))
    y += Inches(0.65)

    # ── Slide 4: 기술 스택 ──
    s4, y = new_content_slide(prs)
    y = section_title(s4, L, y, "기술 스택")
    for group_name, tags in [
        ("Language", ["Swift", "C++", "Objective-C++"]),
        ("Framework / Library", ["SwiftUI", "UIKit", "AppKit", "Combine", "RxSwift", "TCA", "ReactorKit"]),
        ("Apple Platform", ["System Extension", "Network Extension", "App Extension", "XPC", "Daemon", "DeveloperID"]),
        ("Tools / Architecture", ["Clean Architecture", "Tuist", "Swift Package", "Packaging", "RIBs"]),
    ]:
        y = sub_section(s4, L + Inches(0.1), y, group_name, font_size=Pt(11))
        y = skill_tags_row(s4, L + Inches(0.2), y, tags, width=Inches(8.5))
        y += Inches(0.05)


# ═══════════════════════════════════════
#  경력 기술서 슬라이드 생성
# ═══════════════════════════════════════

# -- 원본 데이터 (career.md 기반) --

CAREER_SUMMARY = [
    ("2023.04 - 현재", "프라이빗테크놀로지(주)", "연구소 주임연구원",
     "iOS/macOS 보안 에이전트 개발 (정보보안, 클라이언트)"),
    ("2020.12 - 2022.02", "제이엠트루", "개발연구소 연구원",
     "iOS 앱 개발 및 배포 (앱개발자)"),
]

CAREER_DETAIL = [
    {
        "company": "프라이빗테크놀로지(주)",
        "period": "2023.04 - 현재",
        "role": "연구소 주임연구원 · 개발연구소",
        "duties": [
            "iOS 에이전트 개발 및 App Store, Enterprise 배포",
            "macOS 에이전트, 데몬 개발 및 Installer(package) 배포",
            "제로트러스트 기반 네트워크 보안 솔루션(SASE) 개발",
        ],
        "achievements": [
            "macOS/iOS PacketGO Switch 에이전트 신규 개발 및 출시",
            "UI/UX 성능 향상 - Animation Hitch 약 60% 감소",
            "보안 기능 확인서 iOS, macOS 취득",
            "KISA 실증 완료",
            "CSAP 인증(시큐어 코딩) 대응",
            "TCA, Clean Architecture 도입을 제안하여 팀 개발 프로세스 개선",
        ],
    },
    {
        "company": "제이엠트루",
        "period": "2020.12 - 2022.02 (1년 3개월)",
        "role": "개발연구소 연구원 · 앱개발",
        "duties": [
            "iOS 앱 개발 및 배포",
            "공공데이터 API 연동 및 데이터 처리",
            "앱스토어 배포 및 출시 관리",
        ],
        "achievements": [
            "500명 이상의 사용자 회원가입 결과물 도출",
        ],
    },
]

PROJECTS = [
    {
        "title": "macOS 제로트러스트 기반 SASE 개발",
        "period": "2025.05 - 현재",
        "role": "iOS/macOS 개발 · 프라이빗테크놀로지(주)",
        "overview": "macOS 환경 제로트러스트 기반 SASE(Secure Access Service Edge) 개발",
        "tasks": [
            "강화된 패킷 제어 - 패킷 정보 수집의 변화, 패킷 제어에 대한 로직 변화 개발",
            "테스트 및 유지보수를 위한 모듈화",
            "아키텍처 설계 (TCA, Clean Architecture)",
        ],
        "achievements": None,
        "tags": ["Swift", "C++", "Objective-C++", "SwiftUI", "Combine",
                 "System Extension", "Network Extension", "DeveloperID",
                 "Packaging", "Daemon (privileged)", "Command Line Tool",
                 "XPC", "Swift Package", "Framework", "TCA"],
    },
    {
        "title": "PacketGO Switch 고도화 및 유지보수, POC",
        "period": "2024.05 - 현재",
        "role": "iOS/macOS 개발 · 프라이빗테크놀로지(주)",
        "overview": "제로트러스트 기반 네트워크 보안 에이전트 PacketGO Switch의 기능 고도화 및 고객사 POC 대응",
        "tasks": [
            "사용자 통제, 보안 환경 기능 고도화 (IPv6, 스크린 세이버, 방화벽 등)",
            "UI/UX 성능 향상 - 많은 양의 패킷 제어 시 기록된 정보 업데이트에 대한 UX 성능 저하 해결 (멀티 스레딩, UI 컴포넌트 변경)",
            "고객사 POC 대응",
        ],
        "achievements": ["Animation Hitch 약 60% 감소"],
        "tags": None,
    },
    {
        "title": "PacketGO Switch 보안 기능 확인서 취득",
        "period": "2024.01 - 2024.05",
        "role": "iOS/macOS 개발 · 프라이빗테크놀로지(주)",
        "overview": "PacketGO Switch 보안 기능 확인서 iOS, macOS 취득 프로젝트",
        "tasks": [
            "보안 기능 확인서 iOS, macOS 타겟 OS에 맞는 에이전트 개발",
            "타겟 OS 버전이 제품 OS 버전보다 낮은 이유로 프레임워크 변경 및 개발 (SwiftUI → UIKit, AppKit)",
        ],
        "achievements": None,
        "tags": None,
    },
    {
        "title": "macOS PacketGO Switch 에이전트 개발",
        "period": "2023.06 - 2023.11",
        "role": "iOS/macOS 개발 · 프라이빗테크놀로지(주)",
        "overview": "macOS 환경 제로트러스트 기반 네트워크 보안 에이전트 PacketGO Switch 신규 개발",
        "tasks": [
            "iOS, macOS 통합 UI 개발을 위한 SwiftUI 기반 개발",
            "패킷 제어 모듈(C++ 기반) 개발, Swift 호환성 설계 (Swift ↔ Objc-C++ ↔ C++)",
            "VPN 모듈 개발",
            "사용자 통제 및 환경 제어",
            "비 샌드박스 패키징 및 배포",
            "KISA 실증",
        ],
        "achievements": None,
        "tags": ["Swift", "C++", "Objective-C++", "SwiftUI", "Combine",
                 "System Extension", "Network Extension", "DeveloperID",
                 "Packaging", "Daemon (privileged)", "Command Line Tool", "XPC"],
    },
    {
        "title": "iOS PacketGO Switch 에이전트 개발",
        "period": "2023.06 - 2023.11",
        "role": "iOS 개발 · 프라이빗테크놀로지(주)",
        "overview": "iOS 환경 제로트러스트 기반 네트워크 보안 에이전트 PacketGO Switch 신규 개발",
        "tasks": [
            "iOS, macOS 통합 UI 개발을 위한 SwiftUI 기반 개발",
            "패킷 제어 모듈(C++ 기반) 개발, Swift 호환성 설계 (Swift ↔ Objc-C++ ↔ C++)",
            "VPN 모듈 개발",
            "KISA 실증",
        ],
        "achievements": None,
        "tags": ["Swift", "C++", "Objective-C++", "SwiftUI", "Combine",
                 "App Extension", "Network Extension"],
    },
    {
        "title": "QR 인증 iOS 에이전트 개발",
        "period": "2023.04 - 2023.10",
        "role": "iOS 개발 · 프라이빗테크놀로지(주)",
        "overview": "타 OS 환경 사용자 인증을 위한 QR 인증 방식을 지원하는 iOS 에이전트 개발 및 배포",
        "tasks": [
            "UI 개발 및 비즈니스 로직 구현",
            "앱스토어 배포 및 유지보수",
            "엔터프라이즈 배포 및 유지보수",
            "CSAP 인증(시큐어 코딩)",
        ],
        "achievements": None,
        "tags": ["Swift", "Combine", "UIKit"],
    },
    {
        "title": "리콜박스 애플리케이션 개발",
        "period": "2021.01 - 2021.11 (약 11개월)",
        "role": "iOS 개발 · 제이엠트루",
        "overview": "공공데이터 API를 사용하여 리콜 제품의 정보를 사용자에게 제공하고, 정보에 대해 사용자들이 커뮤니케이션 할 수 있는 \"리콜박스\" 애플리케이션 개발",
        "tasks": [
            "UI/UX 설계 및 화면 구성",
            "공공데이터 API 연동 및 데이터 처리",
            "앱스토어 배포 및 출시 관리",
        ],
        "achievements": ["500명 이상의 사용자 회원가입 결과물 도출"],
        "tags": ["Swift", "RxSwift", "UIKit"],
    },
]


def build_career_slides(prs, with_title_slide=True):
    """경력 기술서 슬라이드 생성"""

    # ── Slide 1: 타이틀 + 경력 요약 테이블만 ──
    if with_title_slide:
        s = add_blank_slide(prs)
        add_orange_sidebar(s, width=Inches(3.2))
        textbox(s, Inches(0.4), Inches(1.2), Inches(2.6), Inches(0.5), "안광빈",
                font_name=FONT_HEADING, font_size=Pt(30), bold=True, color=WHITE)
        textbox(s, Inches(0.4), Inches(1.8), Inches(2.6), Inches(0.3), "경력 기술서",
                font_name=FONT_BODY, font_size=Pt(14), color=WHITE)
        textbox(s, Inches(0.4), Inches(2.2), Inches(2.6), Inches(0.25), "macOS, iOS Developer",
                font_name=FONT_BODY, font_size=Pt(10), color=RGBColor(0xFF, 0xCC, 0xAA))

        rx = Inches(3.6)
        y = Inches(0.5)
        y = section_title(s, rx, y, "경력 요약", font_size=Pt(20))

        add_table(s, rx, y,
                  ["기간", "회사명", "직책/직급", "담당 업무"],
                  CAREER_SUMMARY,
                  [Inches(1.4), Inches(1.6), Inches(1.2), Inches(1.8)])
    else:
        # 통합본: 경력기술서 시작 표시
        s, y = new_content_slide(prs)
        y = section_title(s, L, y, "경력 요약", font_size=Pt(20))
        add_table(s, L, y,
                  ["기간", "회사명", "직책/직급", "담당 업무"],
                  CAREER_SUMMARY,
                  [Inches(2.0), Inches(2.5), Inches(2.0), Inches(2.7)])

    # ── Slide 2~: 경력 상세 ──
    slide, y = new_content_slide(prs)
    y = section_title(slide, L, y, "경력 상세")

    for i, career in enumerate(CAREER_DETAIL):
        # 이 경력 블록에 필요한 높이 추정
        needed = Inches(0.3 + 0.24 + 0.27) + bullet_height(career["duties"])
        needed += Inches(0.27) + bullet_height(career["achievements"])

        if needs_new_slide(y, needed):
            slide, y = new_content_slide(prs)

        if i > 0:
            y = thin_divider(slide, L, y)

        y = timeline_header(slide, L, y, career["company"], career["period"])
        y = role_text(slide, L + Inches(0.1), y, career["role"])

        y = sub_section(slide, L + Inches(0.1), y, "담당 업무", font_size=Pt(10))
        bh = bullet_list(slide, L + Inches(0.2), y, career["duties"], width=Inches(8.3))
        y += bh + Inches(0.05)

        y = sub_section(slide, L + Inches(0.1), y, "주요 성과", font_size=Pt(10))
        bh = bullet_list(slide, L + Inches(0.2), y, career["achievements"], width=Inches(8.3))
        y += bh + Inches(0.05)

    # ── 프로젝트 이력 ──
    # 새 슬라이드에서 시작
    slide, y = new_content_slide(prs)
    y = section_title(slide, L, y, "프로젝트 이력")

    for i, proj in enumerate(PROJECTS):
        # 이 프로젝트 블록에 필요한 높이 추정
        needed = estimate_project_height(
            proj.get("overview"), proj.get("tasks"),
            proj.get("achievements"), proj.get("tags"))

        if needs_new_slide(y, needed):
            slide, y = new_content_slide(prs)

        if i > 0 and y > Inches(0.5):
            y = thin_divider(slide, L, y)

        y = render_project(slide, y,
                           proj["title"], proj["period"], proj["role"],
                           proj.get("overview"), proj.get("tasks"),
                           proj.get("achievements"), proj.get("tags"))


# ═══════════════════════════════════════
#  메인
# ═══════════════════════════════════════
if __name__ == "__main__":
    # 1) 이력서 단독
    prs1 = Presentation(TEMPLATE_PATH)
    delete_all_slides(prs1)
    build_resume_slides(prs1)
    out1 = os.path.join(OUT_DIR, "이력서.pptx")
    prs1.save(out1)
    print(f"Generated: {out1}")

    # 2) 경력기술서 단독
    prs2 = Presentation(TEMPLATE_PATH)
    delete_all_slides(prs2)
    build_career_slides(prs2, with_title_slide=True)
    out2 = os.path.join(OUT_DIR, "경력기술서.pptx")
    prs2.save(out2)
    print(f"Generated: {out2}")

    # 3) 통합본 (이력서 + 경력기술서)
    prs3 = Presentation(TEMPLATE_PATH)
    delete_all_slides(prs3)
    build_resume_slides(prs3)
    build_career_slides(prs3, with_title_slide=False)
    out3 = os.path.join(OUT_DIR, "이력서_경력기술서.pptx")
    prs3.save(out3)
    print(f"Generated: {out3}")
